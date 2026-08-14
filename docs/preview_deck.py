"""Render the generated deck to PNGs so the layout can be reviewed without PowerPoint.

Approximate renderer: solid fills, rectangles, pictures and wrapped text.
Usage: python docs/preview_deck.py [slide numbers...]
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu

DOCS = Path(__file__).resolve().parent
DECK = DOCS / "AI-Skill-Analyser-Client-Demo.pptx"
OUT = DOCS / "preview"
DPI = 110
FONTS = Path(r"C:\Windows\Fonts")
CACHE: dict[tuple[str, int], ImageFont.FreeTypeFont] = {}


def px(value: Emu | int | None) -> int:
    return round(float(value or 0) / 914400 * DPI)


def pt(value: float) -> float:
    return value * DPI / 72


def font(size_pt: float, bold: bool, italic: bool) -> ImageFont.FreeTypeFont:
    name = "segoeuiz" if (bold and italic) else "segoeuib" if bold else "segoeuii" if italic else "segoeui"
    size = max(6, round(size_pt * DPI / 72))
    key = (name, size)
    if key not in CACHE:
        CACHE[key] = ImageFont.truetype(str(FONTS / f"{name}.ttf"), size)
    return CACHE[key]


def rgb(color, default=(40, 40, 40)):
    try:
        if color is None or color.type is None:
            return default
        value = color.rgb
    except (AttributeError, TypeError, ValueError):
        return default
    return (value[0], value[1], value[2])


def shape_fill(shape):
    try:
        if shape.fill.type == 1:
            return rgb(shape.fill.fore_color, (255, 255, 255))
    except (AttributeError, TypeError, ValueError):
        return None
    return None


def shape_line(shape):
    try:
        if shape.line.fill.type != 1:
            return None, 0
        return rgb(shape.line.color, (200, 200, 200)), max(1, px(shape.line.width))
    except (AttributeError, TypeError, ValueError):
        return None, 0


def draw_shape(draw: ImageDraw.ImageDraw, shape) -> None:
    fill = shape_fill(shape)
    line, width = shape_line(shape)
    if fill is None and line is None:
        return
    x0, y0 = px(shape.left), px(shape.top)
    box = [x0, y0, x0 + px(shape.width), y0 + px(shape.height)]
    if box[2] <= box[0] or box[3] <= box[1]:
        return
    label = str(getattr(shape, "shape_type", "")) + (shape.name or "")
    if "ROUNDED" in label.upper():
        radius = min(16, (box[3] - box[1]) // 2, (box[2] - box[0]) // 2)
        draw.rounded_rectangle(box, radius=radius, fill=fill, outline=line, width=width or 1)
    elif "TRIANGLE" in label.upper():
        draw.polygon([(box[0], box[3]), (box[2], box[3]), ((box[0] + box[2]) / 2, box[1])], fill=fill)
    else:
        draw.rectangle(box, fill=fill, outline=line, width=width or 1)


def layout(draw: ImageDraw.ImageDraw, frame, width: int):
    """Return (lines, total_height) where a line is (segments, height, before, after, align)."""
    lines = []
    for paragraph in frame.paragraphs:
        segments = [
            (
                run.text,
                run.font.size.pt if run.font.size else 12,
                bool(run.font.bold),
                bool(run.font.italic),
                rgb(run.font.color),
            )
            for run in paragraph.runs
            if run.text
        ]
        if not segments:
            continue
        spacing = paragraph.line_spacing or 1.0
        before = paragraph.space_before.pt if paragraph.space_before else 0
        after = paragraph.space_after.pt if paragraph.space_after else 0
        align = paragraph.alignment
        biggest = max(size for _, size, _, _, _ in segments)
        line_h = pt(biggest * spacing * 1.24)
        current: list[tuple[str, ImageFont.FreeTypeFont, tuple[int, int, int]]] = []
        used = 0.0
        started = len(lines)
        for text, size, bold, italic, colour in segments:
            glyph = font(size, bold, italic)
            words = text.split(" ")
            for index, word in enumerate(words):
                if not word and index < len(words) - 1:
                    continue
                token = word if not current else " " + word
                advance = draw.textlength(token, font=glyph)
                if used + advance > width and current:
                    lines.append([current, line_h, 0.0, 0.0, align])
                    current = [(word, glyph, colour)]
                    used = draw.textlength(word, font=glyph)
                else:
                    current.append((token, glyph, colour))
                    used += advance
        if current:
            lines.append([current, line_h, 0.0, 0.0, align])
        if len(lines) > started:
            lines[started][2] = pt(before)
            lines[-1][3] = pt(after)
    total = sum(line[1] + line[2] + line[3] for line in lines)
    return lines, total


def draw_text(draw: ImageDraw.ImageDraw, shape) -> list[str]:
    frame = shape.text_frame
    if not frame.text.strip():
        return []
    left = px(shape.left) + px(frame.margin_left)
    top = px(shape.top) + px(frame.margin_top)
    width = px(shape.width) - px(frame.margin_left) - px(frame.margin_right)
    height = px(shape.height) - px(frame.margin_top) - px(frame.margin_bottom)
    if width <= 0:
        return []

    lines, total = layout(draw, frame, width)
    y = float(top)
    if frame.vertical_anchor == MSO_ANCHOR.MIDDLE:
        y += max(0.0, (height - total) / 2)
    elif frame.vertical_anchor == MSO_ANCHOR.BOTTOM:
        y += max(0.0, height - total)

    for segments, line_h, before, after, align in lines:
        y += before
        line_width = sum(draw.textlength(text, font=glyph) for text, glyph, _ in segments)
        x = float(left)
        if align == PP_ALIGN.CENTER:
            x += (width - line_width) / 2
        elif align == PP_ALIGN.RIGHT:
            x += width - line_width
        for text, glyph, colour in segments:
            draw.text((x, y), text, font=glyph, fill=colour)
            x += draw.textlength(text, font=glyph)
        y += line_h + after

    warnings = []
    if total > height + 8:
        warnings.append(f"+{total - height:.0f}px '{frame.text[:44].strip()}'")
    return warnings


def render(indices: list[int] | None = None) -> None:
    prs = Presentation(str(DECK))
    OUT.mkdir(exist_ok=True)
    width, height = px(prs.slide_width), px(prs.slide_height)
    for number, slide in enumerate(prs.slides, start=1):
        if indices and number not in indices:
            continue
        canvas = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(canvas)
        warnings: list[str] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                with Image.open(io.BytesIO(shape.image.blob)) as picture:
                    resized = picture.convert("RGB").resize(
                        (max(1, px(shape.width)), max(1, px(shape.height)))
                    )
                    canvas.paste(resized, (px(shape.left), px(shape.top)))
                continue
            draw_shape(draw, shape)
            if shape.has_text_frame:
                warnings.extend(draw_text(draw, shape))
        canvas.save(OUT / f"slide-{number:02d}.png")
        if warnings:
            print(f"slide {number:02d}: " + " | ".join(warnings))
    print(f"rendered to {OUT}")


if __name__ == "__main__":
    wanted = [int(value) for value in sys.argv[1:]]
    render(wanted or None)
