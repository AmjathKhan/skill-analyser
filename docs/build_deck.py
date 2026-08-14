"""Build the AI Skill Analyser client demo deck.

Usage:  python docs/build_deck.py
Output: docs/AI-Skill-Analyser-Client-Demo.pptx
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

DOCS = Path(__file__).resolve().parent
SHOTS = DOCS / "screenshots"
OUTPUT = DOCS / "AI-Skill-Analyser-Client-Demo.pptx"

INK = RGBColor(0x0F, 0x1E, 0x36)
NAVY = RGBColor(0x0A, 0x18, 0x2E)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
VIOLET = RGBColor(0x7C, 0x4D, 0xFF)
GREEN = RGBColor(0x12, 0xB7, 0x6A)
AMBER = RGBColor(0xF7, 0x90, 0x09)
RED = RGBColor(0xF0, 0x44, 0x38)
TEAL = RGBColor(0x0E, 0x9F, 0xA8)
SLATE = RGBColor(0x5B, 0x6B, 0x84)
MIST = RGBColor(0x8C, 0x9C, 0xB5)
LIGHT = RGBColor(0xF4, 0xF7, 0xFC)
BORDER = RGBColor(0xE1, 0xE8, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.75
BODY_W = SLIDE_W - 2 * MARGIN


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.count = 0

    # ---------- primitives -------------------------------------------------

    def _slide(self, background: RGBColor = WHITE):
        slide = self.prs.slides.add_slide(self.blank)
        self.count += 1
        back = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        back.fill.solid()
        back.fill.fore_color.rgb = background
        back.line.fill.background()
        back.shadow.inherit = False
        return slide

    def rect(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        fill: RGBColor | None = None,
        line: RGBColor | None = None,
        radius: float | None = None,
        line_width: float = 1.0,
    ):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        if radius:
            shape.adjustments[0] = radius
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
            shape.line.width = Pt(line_width)
        shape.shadow.inherit = False
        shape.text_frame.word_wrap = True
        return shape

    def text(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        anchor=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = anchor
        return frame

    def line(
        self,
        frame,
        text: str,
        size: float = 12,
        color: RGBColor = SLATE,
        bold: bool = False,
        space_before: float = 0,
        space_after: float = 0,
        align=PP_ALIGN.LEFT,
        spacing: float = 1.0,
        italic: bool = False,
        first: bool = False,
    ):
        para = frame.paragraphs[0] if first else frame.add_paragraph()
        para.alignment = align
        para.space_before = Pt(space_before)
        para.space_after = Pt(space_after)
        para.line_spacing = spacing
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = FONT
        return para

    def notes(self, slide, text: str) -> None:
        slide.notes_slide.notes_text_frame.text = text

    # ---------- composite pieces ------------------------------------------

    def header(
        self,
        slide,
        kicker: str,
        title: str,
        subtitle: str = "",
        accent: RGBColor = BLUE,
    ) -> None:
        frame = self.text(slide, MARGIN, 0.44, BODY_W, 0.3)
        self.line(frame, kicker.upper(), 10.5, accent, bold=True, first=True)
        frame = self.text(slide, MARGIN, 0.76, BODY_W, 0.62)
        self.line(frame, title, 29, INK, bold=True, first=True)
        if subtitle:
            frame = self.text(slide, MARGIN, 1.42, BODY_W - 0.4, 0.5)
            self.line(frame, subtitle, 13, SLATE, first=True, spacing=1.15)
        self.rect(slide, MARGIN - 0.22, 0.72, 0.055, 0.48, fill=accent)

    def footer(self, slide, dark: bool = False) -> None:
        colour = MIST if dark else RGBColor(0x9A, 0xA8, 0xBC)
        if not dark:
            self.rect(slide, MARGIN, 7.02, BODY_W, 0.012, fill=BORDER)
        frame = self.text(slide, MARGIN, 7.12, 7.0, 0.25)
        self.line(
            frame,
            "AI Skill Analyser  ·  Graph RAG HR Recruitment Platform",
            9,
            colour,
            first=True,
        )
        frame = self.text(slide, SLIDE_W - MARGIN - 1.2, 7.12, 1.2, 0.25)
        self.line(frame, f"{self.count:02d}", 9, colour, align=PP_ALIGN.RIGHT, first=True)

    def card(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        body: str = "",
        accent: RGBColor = BLUE,
        fill: RGBColor = WHITE,
        title_size: float = 13,
        body_size: float = 10.5,
        eyebrow: str = "",
    ):
        self.rect(slide, x, y, w, h, fill=fill, line=BORDER, radius=0.08)
        self.rect(slide, x, y, 0.05, h, fill=accent)
        frame = self.text(slide, x + 0.28, y + 0.22, w - 0.52, h - 0.4)
        first = True
        if eyebrow:
            self.line(frame, eyebrow.upper(), 9, accent, bold=True, first=True, space_after=5)
            first = False
        self.line(frame, title, title_size, INK, bold=True, first=first, spacing=1.05, space_after=7)
        if body:
            self.line(frame, body, body_size, SLATE, spacing=1.22)

    def stat(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        value: str,
        label: str,
        accent: RGBColor = BLUE,
        note: str = "",
        fill: RGBColor = LIGHT,
    ) -> None:
        self.rect(slide, x, y, w, h, fill=fill, line=BORDER, radius=0.1)
        frame = self.text(slide, x + 0.26, y + 0.24, w - 0.5, 0.55)
        self.line(frame, value, 27, accent, bold=True, first=True)
        frame = self.text(slide, x + 0.26, y + 0.82, w - 0.5, 0.6)
        self.line(frame, label, 10.5, INK, bold=True, first=True, spacing=1.15)
        if note:
            self.line(frame, note, 9.5, SLATE, space_before=3, spacing=1.15)

    def chip(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        text: str,
        colour: RGBColor,
        h: float = 0.34,
        size: float = 10,
    ) -> None:
        shape = self.rect(slide, x, y, w, h, fill=colour, radius=0.5)
        frame = shape.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.margin_left = frame.margin_right = 0
        self.line(frame, text, size, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)

    def bullets(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        items: list[str],
        size: float = 11.5,
        gap: float = 9,
        colour: RGBColor = SLATE,
        marker: str = "—",
        marker_colour: RGBColor = BLUE,
        h: float = 4.0,
    ) -> None:
        frame = self.text(slide, x, y, w, h)
        for index, item in enumerate(items):
            para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            para.space_after = Pt(gap)
            para.line_spacing = 1.2
            head = para.add_run()
            head.text = f"{marker}  "
            head.font.size = Pt(size)
            head.font.bold = True
            head.font.color.rgb = marker_colour
            head.font.name = FONT
            if "  ·  " in item:
                lead, rest = item.split("  ·  ", 1)
                strong = para.add_run()
                strong.text = f"{lead}  "
                strong.font.size = Pt(size)
                strong.font.bold = True
                strong.font.color.rgb = INK
                strong.font.name = FONT
                item = rest
            body = para.add_run()
            body.text = item
            body.font.size = Pt(size)
            body.font.color.rgb = colour
            body.font.name = FONT

    def picture(
        self,
        slide,
        name: str,
        x: float,
        y: float,
        w: float,
        caption: str = "",
    ) -> float:
        path = SHOTS / name
        if not path.exists():
            self.rect(slide, x, y, w, w / 2, fill=LIGHT, line=BORDER, radius=0.05)
            return y + w / 2
        with Image.open(path) as image:
            ratio = image.height / image.width
        h = w * ratio
        self.rect(slide, x - 0.045, y - 0.045, w + 0.09, h + 0.09, fill=BORDER, radius=0.03)
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
        if caption:
            frame = self.text(slide, x, y + h + 0.14, w, 0.3)
            self.line(frame, caption, 9.5, MIST, first=True, italic=True)
        return y + h

    def arrow_flow(
        self,
        slide,
        y: float,
        steps: list[tuple[str, str, RGBColor]],
        height: float = 1.5,
    ) -> None:
        count = len(steps)
        gap = 0.16
        width = (BODY_W - gap * (count - 1)) / count
        for index, (title, body, colour) in enumerate(steps):
            x = MARGIN + index * (width + gap)
            self.rect(slide, x, y, width, height, fill=WHITE, line=BORDER, radius=0.09)
            self.rect(slide, x, y, width, 0.055, fill=colour)
            frame = self.text(slide, x + 0.18, y + 0.22, width - 0.36, 0.22)
            self.line(frame, f"STEP {index + 1}", 8.5, colour, bold=True, first=True)
            frame = self.text(slide, x + 0.18, y + 0.46, width - 0.36, 0.28)
            self.line(frame, title, 12, INK, bold=True, first=True)
            frame = self.text(slide, x + 0.18, y + 0.76, width - 0.36, height - 0.9)
            self.line(frame, body, 9.5, SLATE, first=True, spacing=1.2)
            if index < count - 1:
                chevron = slide.shapes.add_shape(
                    MSO_SHAPE.ISOSCELES_TRIANGLE,
                    Inches(x + width + 0.02),
                    Inches(y + height / 2 - 0.07),
                    Inches(0.12),
                    Inches(0.14),
                )
                chevron.rotation = 90
                chevron.fill.solid()
                chevron.fill.fore_color.rgb = BORDER
                chevron.line.fill.background()
                chevron.shadow.inherit = False

    # ---------- slide types -------------------------------------------------

    def content_slide(
        self, kicker: str, title: str, subtitle: str = "", accent: RGBColor = BLUE
    ):
        slide = self._slide()
        self.header(slide, kicker, title, subtitle, accent)
        self.footer(slide)
        return slide

    def section_slide(self, number: str, title: str, subtitle: str):
        slide = self._slide(NAVY)
        self.rect(slide, 0, 0, 0.22, SLIDE_H, fill=BLUE)
        frame = self.text(slide, 1.4, 2.5, 10.5, 0.4)
        self.line(frame, number, 12, VIOLET, bold=True, first=True)
        frame = self.text(slide, 1.4, 2.9, 10.5, 0.9)
        self.line(frame, title, 40, WHITE, bold=True, first=True)
        frame = self.text(slide, 1.4, 3.95, 9.0, 0.6)
        self.line(frame, subtitle, 14, MIST, first=True, spacing=1.2)
        self.footer(slide, dark=True)
        return slide

    def shot_slide(
        self,
        kicker: str,
        title: str,
        subtitle: str,
        image: str,
        points: list[str],
        caption: str = "",
        accent: RGBColor = BLUE,
    ):
        slide = self.content_slide(kicker, title, subtitle, accent)
        bottom = self.picture(slide, image, MARGIN, 2.15, 8.05, caption)
        panel_x = MARGIN + 8.05 + 0.42
        panel_w = SLIDE_W - MARGIN - panel_x
        self.rect(slide, panel_x, 2.15, panel_w, bottom - 2.15, fill=LIGHT, line=BORDER, radius=0.07)
        frame = self.text(slide, panel_x + 0.26, 2.42, panel_w - 0.52, 0.25)
        self.line(frame, "IN THIS VIEW", 9, accent, bold=True, first=True)
        self.bullets(
            slide,
            panel_x + 0.26,
            2.78,
            panel_w - 0.52,
            points,
            size=10.5,
            gap=10,
            marker="›",
            marker_colour=accent,
            h=3.4,
        )
        return slide

    def save(self) -> None:
        self.prs.save(OUTPUT)


def build() -> None:
    deck = Deck()

    # 01 — title -------------------------------------------------------------
    slide = deck._slide(NAVY)
    deck.rect(slide, 0, 0, SLIDE_W, 0.14, fill=BLUE)
    deck.rect(slide, 8.9, 0.14, 4.44, SLIDE_H - 0.14, fill=RGBColor(0x0E, 0x22, 0x3E))
    deck.rect(slide, 8.9, 0.14, 0.02, SLIDE_H - 0.14, fill=RGBColor(0x1B, 0x33, 0x55))
    frame = deck.text(slide, 1.1, 1.62, 7.4, 0.3)
    deck.line(frame, "CLIENT DEMO  ·  AUGUST 2026", 11, VIOLET, bold=True, first=True)
    frame = deck.text(slide, 1.1, 2.05, 7.6, 1.4)
    deck.line(frame, "AI Skill Analyser", 52, WHITE, bold=True, first=True)
    frame = deck.text(slide, 1.1, 3.15, 7.4, 0.6)
    deck.line(frame, "Graph RAG–Based HR Recruitment Platform", 19, RGBColor(0x9F, 0xC4, 0xFF), first=True)
    deck.rect(slide, 1.1, 3.95, 1.3, 0.035, fill=VIOLET)
    frame = deck.text(slide, 1.1, 4.25, 7.0, 0.9)
    deck.line(
        frame,
        "Hire on evidence, not keywords. Upload a resume, get a ranked, "
        "explainable shortlist in seconds.",
        14,
        MIST,
        first=True,
        spacing=1.3,
    )
    for index, label in enumerate(
        ["Explainable AI matching", "Knowledge graph", "Enterprise ready"]
    ):
        deck.chip(slide, 1.1 + index * 2.35, 5.35, 2.15, label, RGBColor(0x18, 0x2F, 0x52), size=9.5)
    for index, (value, label) in enumerate(
        [("145", "skills in the KB"), ("356", "graph nodes"), ("1 351", "relationships"), ("138", "automated tests")]
    ):
        top = 1.75 + index * 1.05
        frame = deck.text(slide, 9.5, top, 3.3, 0.4)
        deck.line(frame, value, 24, WHITE, bold=True, first=True)
        frame = deck.text(slide, 9.5, top + 0.42, 3.3, 0.3)
        deck.line(frame, label, 10, MIST, first=True)
    deck.footer(slide, dark=True)
    deck.notes(
        slide,
        "Open with the promise: every resume becomes structured, comparable, explainable evidence. "
        "The numbers on the right are from the live demo dataset shipped with the platform.",
    )

    # 02 — agenda ------------------------------------------------------------
    slide = deck.content_slide(
        "Agenda",
        "What we will cover in the next 30 minutes",
        "Five parts: the problem we set out to solve, the goal, how the platform works, a live walkthrough, and the impact.",
    )
    agenda = [
        ("01", "The problem", "Where hiring loses time, money and good candidates", RED),
        ("02", "Goal & end state", "What we committed to build and what success looks like", AMBER),
        ("03", "The solution", "Graph RAG architecture, scoring model, security", BLUE),
        ("04", "Live walkthrough", "The working product, screen by screen", VIOLET),
        ("05", "Impact & next steps", "Value, rollout plan and roadmap", GREEN),
    ]
    for index, (number, title, body, colour) in enumerate(agenda):
        y = 2.15 + index * 0.92
        deck.rect(slide, MARGIN, y, BODY_W, 0.8, fill=LIGHT, line=BORDER, radius=0.08)
        deck.rect(slide, MARGIN, y, 0.05, 0.8, fill=colour)
        frame = deck.text(slide, MARGIN + 0.32, y + 0.24, 0.7, 0.3)
        deck.line(frame, number, 15, colour, bold=True, first=True)
        frame = deck.text(slide, MARGIN + 1.05, y + 0.22, 3.4, 0.35)
        deck.line(frame, title, 14, INK, bold=True, first=True)
        frame = deck.text(slide, MARGIN + 4.55, y + 0.25, BODY_W - 4.9, 0.35)
        deck.line(frame, body, 11.5, SLATE, first=True)
    deck.notes(slide, "Set expectations: roughly 10 minutes of context, 15 minutes of live product, 5 minutes of next steps.")

    # 03 — hiring reality ----------------------------------------------------
    slide = deck.content_slide(
        "The problem  ·  01",
        "Where the hiring funnel actually loses time",
        "An illustrative funnel for a single technical role. Replace the volumes with your own and the shape stays the same.",
        RED,
    )
    funnel = [
        ("250", "resumes received", "per open technical role", BLUE),
        ("12.5 hrs", "manual reading", "at ~3 minutes per resume", AMBER),
        ("40", "shortlisted", "ranked by whoever read them", VIOLET),
        ("12", "interviewed", "after a second manual pass", TEAL),
        ("1", "hired", "4–8 weeks after the first resume arrived", GREEN),
    ]
    width = (BODY_W - 0.18 * 4) / 5
    for index, (value, label, note, colour) in enumerate(funnel):
        deck.stat(
            slide,
            MARGIN + index * (width + 0.18),
            2.2,
            width,
            1.62,
            value,
            label,
            colour,
            note,
        )
    deck.rect(slide, MARGIN, 4.15, BODY_W, 1.95, fill=LIGHT, line=BORDER, radius=0.08)
    frame = deck.text(slide, MARGIN + 0.35, 4.42, 5.4, 0.3)
    deck.line(frame, "THE PART NOBODY BUDGETS FOR", 9.5, RED, bold=True, first=True)
    deck.bullets(
        slide,
        MARGIN + 0.35,
        4.78,
        5.3,
        [
            "12.5 hours of reading  ·  before a single conversation happens",
            "20 roles a quarter  ·  is roughly 250 hours, about 1.5 FTE-months of pure screening",
            "Every re-read starts from zero  ·  nothing from the last search is reusable",
        ],
        size=10.5,
        gap=8,
        marker_colour=RED,
        h=1.3,
    )
    frame = deck.text(slide, 6.6, 4.42, 5.9, 0.3)
    deck.line(frame, "AND THE QUALITY COST", 9.5, AMBER, bold=True, first=True)
    deck.bullets(
        slide,
        6.6,
        4.78,
        5.75,
        [
            "Fatigue sets in  ·  resume 200 is not read like resume 2",
            "Two recruiters, two shortlists  ·  ranking is not reproducible",
            "The strongest candidate may simply never be opened",
        ],
        size=10.5,
        gap=8,
        marker_colour=AMBER,
        h=1.3,
    )
    frame = deck.text(slide, MARGIN, 6.28, BODY_W, 0.3)
    deck.line(
        frame,
        "Illustrative model — the demo uses your volumes during the workshop.",
        9.5,
        MIST,
        italic=True,
        first=True,
    )
    deck.notes(
        slide,
        "Ask the client for their real numbers here: resumes per role, roles per quarter, minutes per resume. "
        "The arithmetic is deliberately simple so they can challenge it.",
    )

    # 04 — pain points -------------------------------------------------------
    slide = deck.content_slide(
        "The problem  ·  02",
        "Six pain points we heard from HR and hiring managers",
        "Every one of these is a direct requirement in the platform we built.",
        RED,
    )
    pains = [
        ("Manual screening does not scale", "Headcount grows linearly with applications. Hiring freezes when recruiters go on leave.", RED),
        ("The same skill, ten spellings", "K8s, Kubernetes, kubernetes, K8's, EKS — a keyword filter treats them as five different things.", AMBER),
        ("Keyword ATS filters strings, not meaning", "A Django expert is invisible to a 'Python' filter if the word Python never appears.", VIOLET),
        ("Scoring is inconsistent", "Every recruiter weighs experience, certifications and projects differently, and never the same way twice.", BLUE),
        ("Good candidates stay buried", "Past applicants are never revisited. The same role is sourced from scratch every time.", TEAL),
        ("No defensible audit trail", "'Why was this candidate rejected?' has no answer that survives a compliance review.", GREEN),
    ]
    card_w = (BODY_W - 0.3 * 2) / 3
    for index, (title, body, colour) in enumerate(pains):
        row, col = divmod(index, 3)
        deck.card(
            slide,
            MARGIN + col * (card_w + 0.3),
            2.2 + row * 2.15,
            card_w,
            1.9,
            title,
            body,
            colour,
            eyebrow=f"pain {index + 1:02d}",
        )
    deck.notes(
        slide,
        "Pause here and ask which two hurt the most. Whichever they pick, point forward to the slide that answers it: "
        "normalization for spelling, Graph RAG for meaning, the scoring model for consistency, the audit log for compliance.",
    )

    # 05 — why keyword search fails -----------------------------------------
    slide = deck.content_slide(
        "The problem  ·  03",
        "Why keyword search quietly loses your best candidates",
        "Four failures we can reproduce on demand in any keyword-based ATS.",
        RED,
    )
    deck.rect(slide, MARGIN, 2.15, 7.35, 4.3, fill=WHITE, line=BORDER, radius=0.07)
    frame = deck.text(slide, MARGIN + 0.3, 2.4, 6.8, 0.3)
    deck.line(frame, "WHAT THE FILTER DOES", 9.5, RED, bold=True, first=True)
    failures = [
        ('Search "Kubernetes"', 'Misses K8s, EKS, AKS, "container orchestration" — the candidate is filtered out on spelling.'),
        ('Search "Python"', 'Matches a resume that mentions Python once under Interests, and ranks it beside an eight-year Python engineer.'),
        ('Filter "5+ years"', 'Reads any number near the word "years" instead of computing real, de-overlapped employment intervals.'),
        ("Related skills", 'Django, Flask and Pandas never roll up to Python. AWS never rolls up from EC2, S3 and Lambda.'),
    ]
    for index, (query, effect) in enumerate(failures):
        y = 2.78 + index * 0.9
        deck.rect(slide, MARGIN + 0.3, y, 1.95, 0.34, fill=RGBColor(0xFE, 0xF0, 0xEE), radius=0.3)
        frame = deck.text(slide, MARGIN + 0.3, y + 0.06, 1.95, 0.25)
        deck.line(frame, query, 9.5, RED, bold=True, align=PP_ALIGN.CENTER, first=True)
        frame = deck.text(slide, MARGIN + 2.42, y + 0.02, 4.7, 0.8)
        deck.line(frame, effect, 10.5, SLATE, first=True, spacing=1.2)
    deck.rect(slide, 8.55, 2.15, 4.03, 4.3, fill=NAVY, radius=0.07)
    frame = deck.text(slide, 8.85, 2.42, 3.5, 0.3)
    deck.line(frame, "WHAT IS ACTUALLY NEEDED", 9.5, RGBColor(0x7A, 0xE2, 0xB8), bold=True, first=True)
    needs = [
        ("One canonical skill", "every synonym resolves to the same node"),
        ("A hierarchy", "Django is a child of Python, EC2 of AWS"),
        ("Evidence", "the sentence in the resume that proves the skill"),
        ("A weighted score", "skills, semantics, experience, certs, projects"),
        ("An explanation", "the graph path that produced the ranking"),
    ]
    for index, (title, body) in enumerate(needs):
        y = 2.85 + index * 0.68
        deck.rect(slide, 8.85, y + 0.08, 0.09, 0.09, fill=GREEN, radius=0.5)
        frame = deck.text(slide, 9.1, y, 3.25, 0.28)
        deck.line(frame, title, 11.5, WHITE, bold=True, first=True)
        frame = deck.text(slide, 9.1, y + 0.24, 3.25, 0.3)
        deck.line(frame, body, 9.5, MIST, first=True)
    deck.notes(
        slide,
        "This slide is the bridge into Graph RAG. The right-hand column is literally the feature list of the platform.",
    )

    # 06 — cost of the status quo -------------------------------------------
    slide = deck.content_slide(
        "The problem  ·  04",
        "What the status quo costs the business",
        "Five risks that compound every quarter the process stays manual.",
        RED,
    )
    costs = [
        ("Speed", "Time-to-shortlist", "Days of reading before the first interview. The best candidates accept another offer while the pile is being read.", AMBER),
        ("Money", "Cost per hire", "Screening hours are the largest hidden line item, and they scale with volume rather than with hires.", RED),
        ("Quality", "Wrong shortlists", "Ranking depends on who read the resume and when. Strong profiles are missed, weak ones advance.", VIOLET),
        ("Reuse", "A dead database", "Thousands of parsed resumes sit unused because nothing makes them searchable by capability.", BLUE),
        ("Risk", "No audit trail", "No reproducible reason for a rejection, which is exactly what regulators and candidates ask for.", GREEN),
    ]
    width = (BODY_W - 0.24 * 4) / 5
    for index, (eyebrow, title, body, colour) in enumerate(costs):
        deck.card(
            slide,
            MARGIN + index * (width + 0.24),
            2.2,
            width,
            2.6,
            title,
            body,
            colour,
            eyebrow=eyebrow,
            title_size=13.5,
            body_size=10,
        )
    deck.rect(slide, MARGIN, 5.15, BODY_W, 1.05, fill=NAVY, radius=0.08)
    frame = deck.text(slide, MARGIN + 0.4, 5.42, 11.0, 0.55)
    deck.line(
        frame,
        "The bottleneck is not a lack of candidates. It is the absence of a machine-readable, "
        "comparable representation of what each candidate can actually do.",
        14,
        WHITE,
        bold=True,
        first=True,
        spacing=1.2,
    )
    deck.notes(slide, "Land the reframe: this is a data structure problem, not a headcount problem. That is what the knowledge graph fixes.")

    # 07 — the goal ----------------------------------------------------------
    slide = deck.content_slide(
        "The goal  ·  05",
        "What we set out to build",
        "One objective, five measurable commitments.",
        AMBER,
    )
    deck.rect(slide, MARGIN, 2.15, BODY_W, 1.12, fill=RGBColor(0xFF, 0xF8, 0xEC), line=RGBColor(0xF7, 0xD9, 0xA8), radius=0.08)
    frame = deck.text(slide, MARGIN + 0.42, 2.36, 11.0, 0.3)
    deck.line(frame, "PROJECT GOAL", 9.5, AMBER, bold=True, first=True)
    frame = deck.text(slide, MARGIN + 0.42, 2.64, 11.2, 0.6)
    deck.line(
        frame,
        "Turn every resume into structured, comparable, explainable evidence — and rank candidates "
        "against any role in seconds.",
        16,
        INK,
        bold=True,
        first=True,
        spacing=1.15,
    )
    goals = [
        ("G1", "Ingest anything", "PDF, DOC, DOCX and TXT at batch scale, with OCR fallback for scanned resumes and checksum-based duplicate detection.", BLUE),
        ("G2", "One shared vocabulary", "Every extracted phrase normalized against an authoritative Skills Knowledge Base you own and can replace.", VIOLET),
        ("G3", "Connect the evidence", "Candidates, skills, technologies, companies, certifications, projects and job roles in one traversable graph.", TEAL),
        ("G4", "Rank, and justify it", "A weighted score you can tune per role, with the component breakdown and graph path behind every number.", GREEN),
        ("G5", "Make it usable", "Plain-English search, dashboards, exportable reports and role-based access for the whole hiring team.", AMBER),
    ]
    card_w = (BODY_W - 0.22 * 2) / 3
    for index, (tag, title, body, colour) in enumerate(goals):
        row, col = divmod(index, 3)
        deck.card(
            slide,
            MARGIN + col * (card_w + 0.22),
            3.52 + row * 1.72,
            card_w,
            1.52,
            title,
            body,
            colour,
            eyebrow=tag,
            body_size=10,
        )
    slot_x = MARGIN + 2 * (card_w + 0.22)
    deck.rect(slide, slot_x, 5.24, card_w, 1.52, fill=NAVY, radius=0.08)
    frame = deck.text(slide, slot_x + 0.3, 5.5, card_w - 0.6, 0.28)
    deck.line(frame, "STATUS", 9, RGBColor(0x9F, 0xC4, 0xFF), bold=True, first=True)
    frame = deck.text(slide, slot_x + 0.3, 5.8, card_w - 0.6, 0.8)
    deck.line(
        frame,
        "All five are built, tested and demonstrated in the walkthrough that follows.",
        12,
        WHITE,
        bold=True,
        first=True,
        spacing=1.2,
    )
    deck.notes(slide, "These five goals map one-to-one onto the five product areas shown in the walkthrough.")

    # 08 — the end goal ------------------------------------------------------
    slide = deck.content_slide(
        "The end goal  ·  06",
        "What success looks like once this is live",
        "The target state we are demonstrating today, and the metrics that prove it.",
        GREEN,
    )
    outcomes = [
        ("Every resume is searchable seconds after upload", "Parsing, normalization, embedding and graph insertion happen in one pipeline — no overnight batch."),
        ("The organisation shares one definition of a skill", "Recruiters, hiring managers and reports all resolve to the same taxonomy node."),
        ("Shortlists take minutes, not days", "Describe the role, run the match, review a ranked list with evidence attached."),
        ("No black-box decisions", "Every score carries its breakdown, matched and missing skills, strengths and gaps."),
        ("The talent graph compounds", "Each resume makes future searches better; past applicants resurface automatically for new roles."),
    ]
    for index, (title, body) in enumerate(outcomes):
        y = 2.2 + index * 0.79
        deck.rect(slide, MARGIN, y, 7.6, 0.68, fill=LIGHT, line=BORDER, radius=0.07)
        tick = deck.rect(slide, MARGIN + 0.22, y + 0.19, 0.3, 0.3, fill=GREEN, radius=0.5)
        frame = tick.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = deck.line(frame, "✓", 11, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
        para.runs[0].font.name = "Segoe UI Symbol"
        frame = deck.text(slide, MARGIN + 0.68, y + 0.12, 6.7, 0.25)
        deck.line(frame, title, 11.5, INK, bold=True, first=True)
        frame = deck.text(slide, MARGIN + 0.68, y + 0.36, 6.7, 0.28)
        deck.line(frame, body, 9.5, SLATE, first=True)
    deck.rect(slide, 8.75, 2.2, 3.83, 3.91, fill=NAVY, radius=0.08)
    frame = deck.text(slide, 9.05, 2.45, 3.3, 0.3)
    deck.line(frame, "NORTH-STAR METRICS", 9.5, RGBColor(0x9F, 0xC4, 0xFF), bold=True, first=True)
    metrics = [
        ("Time to shortlist", "days → minutes"),
        ("Screening hours per role", "12.5 → under 1"),
        ("Candidates actually reviewed", "100% of the pool"),
        ("Ranking consistency", "identical inputs, identical result"),
        ("Decisions with an audit trail", "100%"),
    ]
    for index, (label, value) in enumerate(metrics):
        y = 2.85 + index * 0.63
        frame = deck.text(slide, 9.05, y, 3.3, 0.26)
        deck.line(frame, label, 10, MIST, first=True)
        frame = deck.text(slide, 9.05, y + 0.22, 3.3, 0.28)
        deck.line(frame, value, 12.5, WHITE, bold=True, first=True)
    deck.notes(
        slide,
        "This is the slide to return to at the end. Everything in the walkthrough should be traceable back to one of these outcomes.",
    )

    # 09 — solution at a glance ---------------------------------------------
    slide = deck.content_slide(
        "The solution  ·  07",
        "AI Skill Analyser at a glance",
        "Six capabilities, one pipeline, one graph — delivered as a single deployable platform.",
    )
    tiles = [
        ("Resume intelligence", "Multi-engine parsing with OCR fallback. Personal details, work history, education, projects, certifications and skills extracted with the evidence sentence kept.", BLUE),
        ("Skills Knowledge Base", "145 skills across 57 categories with synonyms, parent/child hierarchy, related technologies and job-role mappings. Owned by you as a CSV.", VIOLET),
        ("Knowledge graph", "9 node types and 12 relationship types linking candidates to everything they have used, held, built and studied.", TEAL),
        ("Explainable matching", "Weighted scoring across skills, semantics, experience, certifications and projects, with the graph path behind every point awarded.", GREEN),
        ("Natural-language search", "Ask in plain English across hybrid, semantic, keyword, graph and skill modes, and get a RAG-generated answer with citations.", AMBER),
        ("Reports & governance", "Dashboards, skill-gap analysis, PDF/CSV/Excel exports, role-based access and a complete audit trail.", RED),
    ]
    card_w = (BODY_W - 0.3 * 2) / 3
    for index, (title, body, colour) in enumerate(tiles):
        row, col = divmod(index, 3)
        deck.card(
            slide,
            MARGIN + col * (card_w + 0.3),
            2.2 + row * 2.15,
            card_w,
            1.92,
            title,
            body,
            colour,
            body_size=10,
        )
    deck.notes(slide, "One line per tile. Do not go deep here — each one gets its own screen in the walkthrough.")

    # 10 — why graph rag -----------------------------------------------------
    slide = deck.content_slide(
        "The solution  ·  08",
        "Why Graph RAG, and not keyword or vector search alone",
        "Vector search understands language. A graph understands structure. Recruiting needs both.",
        VIOLET,
    )
    columns = [
        ("Keyword / ATS", "String matching", RED, ["No", "No", "Partial", "No", "No", "Weak"]),
        ("Vector-only AI", "Embedding similarity", AMBER, ["Yes", "No", "No", "No", "Yes", "Weak"]),
        ("Graph RAG  ·  ours", "Graph + vectors + rules", GREEN, ["Yes", "Yes", "Yes", "Yes", "Yes", "Full"]),
    ]
    rows = [
        "Handles synonyms (K8s = Kubernetes)",
        "Understands hierarchy (Django rolls up to Python)",
        "Keeps evidence for each claim",
        "Explains why a candidate ranked where they did",
        "Handles phrasing never seen before",
        "Auditable for compliance",
    ]
    label_w = 4.9
    col_w = (BODY_W - label_w - 0.3) / 3
    for index, (title, sub, colour, _) in enumerate(columns):
        x = MARGIN + label_w + 0.3 + index * col_w
        deck.rect(slide, x, 2.15, col_w - 0.12, 0.72, fill=colour, radius=0.07)
        frame = deck.text(slide, x + 0.1, 2.28, col_w - 0.32, 0.28)
        deck.line(frame, title, 11.5, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
        frame = deck.text(slide, x + 0.1, 2.54, col_w - 0.32, 0.25)
        deck.line(frame, sub, 9, WHITE, align=PP_ALIGN.CENTER, first=True)
    for row_index, label in enumerate(rows):
        y = 3.0 + row_index * 0.63
        fill = LIGHT if row_index % 2 == 0 else WHITE
        deck.rect(slide, MARGIN, y, BODY_W, 0.56, fill=fill, line=BORDER, radius=0.05)
        frame = deck.text(slide, MARGIN + 0.25, y + 0.16, label_w - 0.3, 0.3)
        deck.line(frame, label, 11, INK, first=True)
        for col_index, (_, _, colour, values) in enumerate(columns):
            x = MARGIN + label_w + 0.3 + col_index * col_w
            value = values[row_index]
            tone = GREEN if value in {"Yes", "Full"} else (AMBER if value == "Partial" else RED)
            frame = deck.text(slide, x, y + 0.16, col_w - 0.12, 0.3)
            deck.line(frame, value, 11, tone, bold=True, align=PP_ALIGN.CENTER, first=True)
    deck.notes(
        slide,
        "The honest version of this slide: vector search alone gets you fuzzy recall but no reasons. "
        "The graph is what makes the answer defensible.",
    )

    # 11 — pipeline ----------------------------------------------------------
    slide = deck.content_slide(
        "The solution  ·  09",
        "How it works: resume to ranked shortlist in six steps",
        "One pipeline runs on upload. Everything downstream — search, matching, reports — reads from its output.",
        TEAL,
    )
    deck.arrow_flow(
        slide,
        2.2,
        [
            ("Ingest", "pdfplumber → PyMuPDF → python-docx → OCR fallback. Checksums catch duplicates.", BLUE),
            ("Extract", "Personal details, work history, education, projects, certifications, skill mentions.", VIOLET),
            ("Normalize", "Every phrase resolved to the Skills KB with match type, confidence and evidence.", TEAL),
            ("Build graph", "Candidate linked to skills, companies, projects, certifications and roles.", GREEN),
            ("Embed", "Profile, resume chunks and skill text vectorised for cosine search.", AMBER),
            ("Retrieve & rank", "Query parsed, expanded through the graph, fused with vectors, scored, explained.", RED),
        ],
        height=1.85,
    )
    deck.rect(slide, MARGIN, 4.45, BODY_W, 1.75, fill=LIGHT, line=BORDER, radius=0.08)
    frame = deck.text(slide, MARGIN + 0.35, 4.7, 11.2, 0.3)
    deck.line(frame, "WHAT MAKES THIS PIPELINE DIFFERENT", 9.5, TEAL, bold=True, first=True)
    deck.bullets(
        slide,
        MARGIN + 0.35,
        5.05,
        5.5,
        [
            "Provenance is kept  ·  every skill carries its match type, confidence and source sentence",
            "Experience is computed  ·  merged, de-overlapped employment intervals, not a regex on 'years'",
        ],
        size=10.5,
        gap=8,
        marker_colour=TEAL,
        h=1.1,
    )
    deck.bullets(
        slide,
        6.5,
        5.05,
        5.9,
        [
            "Re-uploads are safe  ·  checksum duplicate detection instead of a second candidate record",
            "Every AI dependency has an offline fallback  ·  no model downloads, no API keys, no GPU required",
        ],
        size=10.5,
        gap=8,
        marker_colour=TEAL,
        h=1.1,
    )
    deck.notes(slide, "Emphasise step 3. Normalization is the step that makes every later number comparable across candidates.")

    # 12 — skills knowledge base --------------------------------------------
    slide = deck.content_slide(
        "The solution  ·  10",
        "The Skills Knowledge Base is your source of truth",
        "A CSV you own, version and replace. Swap the file, re-import, and the whole platform speaks your vocabulary.",
        VIOLET,
    )
    deck.picture(slide, "asa-19-skills-kb.png", MARGIN, 2.2, 6.4, "Skills page — taxonomy browser with category grouping and CSV re-import")
    deck.rect(slide, 7.5, 2.2, 5.08, 1.5, fill=LIGHT, line=BORDER, radius=0.07)
    frame = deck.text(slide, 7.78, 2.42, 4.6, 0.3)
    deck.line(frame, "ONE ROW PER SKILL", 9.5, VIOLET, bold=True, first=True)
    frame = deck.text(slide, 7.78, 2.72, 4.6, 0.9)
    deck.line(
        frame,
        "skill_id · skill_name · category · parent_skill · related_skills · "
        "technology_stack · job_role · experience_level · skill_synonyms · skill_description",
        10,
        SLATE,
        first=True,
        spacing=1.25,
    )
    frame = deck.text(slide, 7.5, 3.92, 5.0, 0.3)
    deck.line(frame, "HOW A RAW PHRASE IS RESOLVED", 9.5, VIOLET, bold=True, first=True)
    ladder = [
        ("1  Exact", "Node.js → nodejs", GREEN, "confidence 1.00"),
        ("2  Synonym", "K8s → Kubernetes", TEAL, "confidence 0.95"),
        ("3  Fuzzy", "Kubernetess → Kubernetes", AMBER, "confidence ~0.80"),
        ("4  Graph", "Django → Python (parent)", VIOLET, "confidence 0.50–0.85"),
    ]
    for index, (step, example, colour, note) in enumerate(ladder):
        y = 4.28 + index * 0.56
        deck.rect(slide, 7.5, y, 5.08, 0.48, fill=WHITE, line=BORDER, radius=0.06)
        deck.rect(slide, 7.5, y, 0.05, 0.48, fill=colour)
        frame = deck.text(slide, 7.72, y + 0.13, 1.15, 0.25)
        deck.line(frame, step, 10, colour, bold=True, first=True)
        frame = deck.text(slide, 8.9, y + 0.13, 2.35, 0.25)
        deck.line(frame, example, 10, INK, first=True)
        frame = deck.text(slide, 11.25, y + 0.14, 1.2, 0.25)
        deck.line(frame, note, 8.5, MIST, align=PP_ALIGN.RIGHT, first=True)
    deck.notes(
        slide,
        "Key ownership message: the client is not locked into our taxonomy. They upload their own CSV from the Skills page and re-import.",
    )

    # 13 — knowledge graph model --------------------------------------------
    slide = deck.content_slide(
        "The solution  ·  11",
        "The knowledge graph model",
        "Nine node types and twelve relationship types. This is what turns isolated resumes into a connected talent asset.",
        TEAL,
    )
    frame = deck.text(slide, MARGIN, 2.15, 5.9, 0.3)
    deck.line(frame, "NODES", 9.5, TEAL, bold=True, first=True)
    nodes = ["Candidate", "Skill", "Category", "Technology", "JobRole", "Company", "Project", "Certification", "Education"]
    for index, name in enumerate(nodes):
        row, col = divmod(index, 3)
        x = MARGIN + col * 1.95
        y = 2.5 + row * 0.56
        deck.rect(slide, x, y, 1.82, 0.46, fill=LIGHT, line=BORDER, radius=0.3)
        frame = deck.text(slide, x, y + 0.12, 1.82, 0.25)
        deck.line(frame, name, 10.5, INK, bold=True, align=PP_ALIGN.CENTER, first=True)
    frame = deck.text(slide, MARGIN, 4.35, 5.9, 0.3)
    deck.line(frame, "RELATIONSHIPS", 9.5, TEAL, bold=True, first=True)
    edges = [
        "HAS_SKILL", "WORKED_AT", "USED_SKILL", "COMPLETED",
        "HOLDS", "STUDIED_AT", "BELONGS_TO", "PART_OF",
        "RELATED_TO", "PARENT_OF", "DEPENDS_ON", "REQUIRED_FOR",
    ]
    for index, name in enumerate(edges):
        row, col = divmod(index, 3)
        x = MARGIN + col * 1.95
        y = 4.7 + row * 0.5
        frame = deck.text(slide, x, y, 1.9, 0.28)
        deck.line(frame, name, 9.5, VIOLET, bold=True, first=True)
    deck.picture(slide, "asa-10-graph-canvas.png", 6.85, 2.2, 5.73, "Live knowledge graph from the five demo resumes")
    deck.rect(slide, 6.85, 5.25, 5.73, 1.0, fill=NAVY, radius=0.07)
    frame = deck.text(slide, 7.12, 5.45, 5.2, 0.28)
    deck.line(frame, "A PATH THE MATCHER ACTUALLY WALKED", 9, RGBColor(0x9F, 0xC4, 0xFF), bold=True, first=True)
    frame = deck.text(slide, 7.12, 5.72, 5.2, 0.45)
    deck.line(
        frame,
        "AWS —[PARENT_OF]→ EC2 (w=0.7)   ·   AWS —[RELATED_TO]→ Terraform (w=0.6)",
        10,
        WHITE,
        first=True,
        spacing=1.2,
    )
    deck.notes(
        slide,
        "The demo dataset is only five resumes and already produces 356 nodes and 1,351 relationships. "
        "Point out that density is the source of the recall advantage.",
    )

    # 14 — scoring model -----------------------------------------------------
    slide = deck.content_slide(
        "The solution  ·  12",
        "The scoring model, in full",
        "No hidden weights. The formula, the credit ladder and the bands are all visible and tunable per role.",
        GREEN,
    )
    deck.rect(slide, MARGIN, 2.15, BODY_W, 0.95, fill=NAVY, radius=0.08)
    frame = deck.text(slide, MARGIN + 0.4, 2.42, 11.2, 0.45)
    deck.line(
        frame,
        "overall  =  0.40 · skill  +  0.20 · semantic  +  0.20 · experience  "
        "+  0.10 · certification  +  0.10 · project",
        15,
        WHITE,
        bold=True,
        first=True,
    )
    frame = deck.text(slide, MARGIN, 3.28, 5.9, 0.3)
    deck.line(frame, "SKILL CREDIT DEGRADES WITH MATCH QUALITY", 9.5, GREEN, bold=True, first=True)
    ladder = [
        ("Exact match", "1.00", GREEN),
        ("Synonym match", "0.95", TEAL),
        ("Fuzzy match", "~0.80", AMBER),
        ("Related / parent / child", "0.55 – 0.85", VIOLET),
        ("Graph-only inference", "0.50", BLUE),
    ]
    for index, (label, value, colour) in enumerate(ladder):
        y = 3.64 + index * 0.52
        deck.rect(slide, MARGIN, y, 5.85, 0.44, fill=LIGHT, line=BORDER, radius=0.05)
        frame = deck.text(slide, MARGIN + 0.25, y + 0.11, 4.0, 0.25)
        deck.line(frame, label, 11, INK, first=True)
        frame = deck.text(slide, MARGIN + 4.2, y + 0.11, 1.4, 0.25)
        deck.line(frame, value, 11, colour, bold=True, align=PP_ALIGN.RIGHT, first=True)
    frame = deck.text(slide, 7.1, 3.28, 5.5, 0.3)
    deck.line(frame, "RECOMMENDATION BANDS", 9.5, GREEN, bold=True, first=True)
    bands = [
        ("Highly Recommended", "85 – 100", GREEN),
        ("Recommended", "70 – 84", BLUE),
        ("Consider", "55 – 69", AMBER),
        ("Not Recommended", "below 55", RED),
    ]
    for index, (label, value, colour) in enumerate(bands):
        y = 3.64 + index * 0.52
        deck.rect(slide, 7.1, y, 5.48, 0.44, fill=WHITE, line=BORDER, radius=0.05)
        deck.rect(slide, 7.1, y, 0.05, 0.44, fill=colour)
        frame = deck.text(slide, 7.35, y + 0.11, 3.4, 0.25)
        deck.line(frame, label, 11, INK, bold=True, first=True)
        frame = deck.text(slide, 10.8, y + 0.11, 1.6, 0.25)
        deck.line(frame, value, 11, colour, bold=True, align=PP_ALIGN.RIGHT, first=True)
    deck.rect(slide, 7.1, 5.72, 5.48, 0.75, fill=RGBColor(0xFE, 0xF0, 0xEE), line=RGBColor(0xF9, 0xC8, 0xC2), radius=0.06)
    frame = deck.text(slide, 7.35, 5.88, 5.0, 0.5)
    deck.line(
        frame,
        "A missing mandatory skill caps the candidate regardless of everything else.",
        10.5,
        RED,
        bold=True,
        first=True,
        spacing=1.15,
    )
    deck.notes(
        slide,
        "Defaults are 40/20/20/10/10 and every weight is a slider on the Skill Match page, so a role can be re-weighted live in the demo.",
    )

    # 15 — architecture ------------------------------------------------------
    slide = deck.content_slide(
        "The solution  ·  13",
        "System architecture",
        "A conventional, well-separated service architecture — nothing exotic to operate.",
    )
    layers = [
        ("Client", "React 18 SPA · TypeScript · Material UI · TanStack Query · Recharts · force-directed graph canvas", BLUE),
        ("API", "FastAPI · JWT auth with refresh · RBAC dependencies · rate limiting · OpenAPI docs · structured errors", VIOLET),
        ("Services", "Resume processing · candidates · matching · search · graph · reports · dashboards · audit", TEAL),
        ("AI & Graph", "Parsers · entity extraction · skill taxonomy · embeddings · vector store · Graph RAG retriever · reasoning", GREEN),
        ("Data", "PostgreSQL + pgvector · Redis · Neo4j or NetworkX · Celery workers · encrypted file storage", AMBER),
    ]
    for index, (name, body, colour) in enumerate(layers):
        y = 2.2 + index * 0.83
        deck.rect(slide, MARGIN, y, BODY_W, 0.72, fill=WHITE, line=BORDER, radius=0.07)
        deck.rect(slide, MARGIN, y, 1.75, 0.72, fill=colour, radius=0.07)
        deck.rect(slide, MARGIN + 1.6, y, 0.15, 0.72, fill=colour)
        frame = deck.text(slide, MARGIN + 0.2, y + 0.24, 1.4, 0.3)
        deck.line(frame, name, 12, WHITE, bold=True, first=True)
        frame = deck.text(slide, MARGIN + 2.0, y + 0.24, BODY_W - 2.3, 0.3)
        deck.line(frame, body, 11, SLATE, first=True)
    deck.rect(slide, MARGIN, 6.4, BODY_W, 0.5, fill=LIGHT, line=BORDER, radius=0.06)
    frame = deck.text(slide, MARGIN + 0.3, 6.52, 11.5, 0.3)
    deck.line(
        frame,
        "Every heavy dependency degrades gracefully: no FAISS → NumPy, no Neo4j → NetworkX, "
        "no model download → deterministic embeddings, no API key → template reasoning.",
        10.5,
        SLATE,
        first=True,
    )
    deck.notes(slide, "The fallback line matters for procurement: the platform runs air-gapped with no external AI calls if required.")

    # 16 — technology stack --------------------------------------------------
    slide = deck.content_slide(
        "The solution  ·  14",
        "Technology stack",
        "Mainstream, supportable technology your team can hire for and operate.",
    )
    stacks = [
        ("Frontend", BLUE, ["React 18 + TypeScript", "Vite build tooling", "Material UI design system", "TanStack Query", "Recharts + force graph", "Light and dark themes"]),
        ("Backend", VIOLET, ["Python 3.11 + FastAPI", "SQLAlchemy 2 ORM", "Alembic migrations", "Pydantic v2 contracts", "Celery + Redis workers", "Structured logging"]),
        ("AI / ML", TEAL, ["spaCy entity extraction", "Sentence Transformers", "pgvector / FAISS / NumPy", "NetworkX + Neo4j", "Graph RAG retriever", "Template or OpenAI reasoning"]),
        ("Data & Ops", AMBER, ["PostgreSQL 15 + pgvector", "Redis cache and queue", "Neo4j graph database", "Docker Compose", "138 automated tests", "Ruff + MyPy quality gates"]),
    ]
    card_w = (BODY_W - 0.3 * 3) / 4
    for index, (name, colour, items) in enumerate(stacks):
        x = MARGIN + index * (card_w + 0.3)
        deck.rect(slide, x, 2.2, card_w, 4.05, fill=WHITE, line=BORDER, radius=0.08)
        deck.rect(slide, x, 2.2, card_w, 0.62, fill=colour, radius=0.08)
        deck.rect(slide, x, 2.66, card_w, 0.16, fill=colour)
        frame = deck.text(slide, x + 0.25, 2.38, card_w - 0.5, 0.3)
        deck.line(frame, name, 13, WHITE, bold=True, first=True)
        for item_index, item in enumerate(items):
            y = 3.02 + item_index * 0.5
            deck.rect(slide, x + 0.25, y + 0.12, 0.07, 0.07, fill=colour, radius=0.5)
            frame = deck.text(slide, x + 0.45, y, card_w - 0.7, 0.4)
            deck.line(frame, item, 10.5, SLATE, first=True, spacing=1.15)
    deck.notes(slide, "No proprietary lock-in. Everything here is open source and deployable in the client's own cloud.")

    # 17 — security ----------------------------------------------------------
    slide = deck.content_slide(
        "The solution  ·  15",
        "Security, access control and governance",
        "Built for an enterprise HR function handling personal data.",
        RED,
    )
    security = [
        ("Authentication", "JWT access and refresh tokens, bcrypt password hashing, session listing and revocation, password reset flow.", BLUE),
        ("Authorisation", "A permission matrix enforced as API dependencies, so the UI and the API can never drift apart.", VIOLET),
        ("Audit", "Every login, upload, edit, status change and export recorded with actor, target and timestamp.", TEAL),
        ("Data protection", "Optional Fernet encryption of stored resumes, soft delete, checksum de-duplication, configurable retention.", GREEN),
        ("Hardening", "Rate limiting, strict schema validation, structured error handling, API docs disabled in production.", AMBER),
        ("Data residency", "Runs entirely inside your infrastructure. With offline models, no candidate data leaves the network.", RED),
    ]
    card_w = (BODY_W - 0.3 * 2) / 3
    for index, (title, body, colour) in enumerate(security):
        row, col = divmod(index, 3)
        deck.card(
            slide,
            MARGIN + col * (card_w + 0.3),
            2.2 + row * 1.72,
            card_w,
            1.5,
            title,
            body,
            colour,
            body_size=10,
        )
    deck.rect(slide, MARGIN, 5.72, BODY_W, 0.72, fill=LIGHT, line=BORDER, radius=0.07)
    roles = [
        ("HR Admin", "Everything, including users, audit and the Skills KB", BLUE),
        ("Recruiter", "Upload, edit candidates, match, search, report", VIOLET),
        ("Hiring Manager", "Read candidates, match, search, report", TEAL),
    ]
    for index, (name, scope, colour) in enumerate(roles):
        x = MARGIN + 0.3 + index * 3.95
        deck.rect(slide, x, 5.9, 0.09, 0.36, fill=colour, radius=0.3)
        frame = deck.text(slide, x + 0.22, 5.88, 3.5, 0.24)
        deck.line(frame, name, 11, INK, bold=True, first=True)
        frame = deck.text(slide, x + 0.22, 6.1, 3.5, 0.24)
        deck.line(frame, scope, 9, SLATE, first=True)
    deck.notes(slide, "Three seeded roles ship with the platform; the permission matrix is configuration, so more roles are additive.")

    # 18 — section divider ---------------------------------------------------
    slide = deck.section_slide(
        "PART 04",
        "Live product walkthrough",
        "Eleven screens from the running platform — every screenshot on the following slides is the "
        "real application, not a mockup.",
    )
    deck.notes(slide, "Switch to the live app here if the room has network access; the slides are the fallback.")

    # 19 — login -------------------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  01",
        "Secure, role-aware access",
        "JWT authentication with refresh tokens, remember-me sessions and a self-service password reset.",
        "asa-01-login.png",
        [
            "Three seeded roles  ·  HR Admin, Recruiter, Hiring Manager",
            "Navigation and actions adapt to the signed-in user's permissions",
            "Active sessions are visible and revocable from Settings",
            "Every sign-in is written to the audit log",
        ],
        "Sign-in — the platform is permission-aware from the first screen",
    )

    # 20 — dashboard ---------------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  02",
        "Dashboard: the state of the pipeline at a glance",
        "Live counts, parse health, average experience and match quality the moment resumes land.",
        "asa-02-dashboard.png",
        [
            "Candidates, resumes processed, shortlisted and pending review",
            "Average experience and average match score across the pool",
            "Top skills and technologies in the talent pool, charted",
            "Recent activity feed and AI-generated recommendations",
        ],
        "Dashboard — recruitment overview with live KPIs",
    )

    # 21 — upload ------------------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  03",
        "Upload: drag, drop, done",
        "Batch upload up to 25 resumes at a time. Parsing, normalization, embedding and graph insertion run automatically.",
        "asa-04-upload.png",
        [
            "PDF, DOC, DOCX and TXT  ·  OCR fallback for scanned documents",
            "Per-file result: extracted name, skills found, parse time, warnings",
            "Duplicate detection by checksum — no accidental double records",
            "Runs inline or on Celery workers for large batches",
        ],
        "Upload — batch ingestion with per-file parse feedback",
        VIOLET,
    )

    # 22 — candidates --------------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  04",
        "Candidates: one searchable, filterable talent pool",
        "Server-side pagination over the whole database, with the filters recruiters actually ask for.",
        "asa-06-candidates-table.png",
        [
            "Filter by skill, status, minimum experience, location and company",
            "Top skills and AI score shown inline for fast triage",
            "Status workflow: pending review, shortlisted, rejected, hired",
            "One-click CSV export of the filtered set",
        ],
        "Candidate list — filters, inline skills and AI score",
        VIOLET,
    )

    # 23 — profile -----------------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  05",
        "Candidate profile: normalized skills with evidence",
        "Raw resume phrases resolved into the taxonomy, grouped by category, each one traceable back to the source text.",
        "asa-08-profile-skills.png",
        [
            "Skills grouped by category, not a flat keyword cloud",
            "Match type and confidence retained for every skill",
            "Timeline, education, projects, certifications and recruiter notes",
            "Similar candidates and the candidate's own subgraph, one tab away",
        ],
        "Candidate profile — skills normalized against the Knowledge Base",
        TEAL,
    )

    # 24 — skill match -------------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  06",
        "AI skill match: a ranked shortlist in under a second",
        "Describe the role, run the match, and get every candidate scored, banded and ordered.",
        "asa-11-match-results.png",
        [
            "Required, mandatory and preferred skills, plus certifications",
            "Live weight sliders — re-rank for a role in real time",
            "Recommendation bands: Highly Recommended → Not Recommended",
            "This run evaluated the pool in 98 ms",
        ],
        "Skill Match — ranked candidates with matched skills and score bands",
        GREEN,
    )

    # 25 — score breakdown ---------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  07",
        "Every score decomposes into five components",
        "Side-by-side comparison shows exactly where each candidate gained and lost points.",
        "asa-13-explainability.png",
        [
            "Skill, semantic, experience, certification and project shown separately",
            "Compare the top candidates on the same axes, at a glance",
            "Change a weight and the whole comparison recalculates",
            "The same numbers are available through the API for your own reporting",
        ],
        "Score comparison — component breakdown per candidate",
        GREEN,
    )

    # 26 — explainability ----------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  08",
        "\"Why this ranking?\" — answered on every candidate",
        "A written justification, the strengths, the gaps and interview questions generated from the evidence.",
        "asa-14-explanation.png",
        [
            "Plain-language explanation of how the score was reached",
            "Strengths with proficiency, and gaps stated explicitly",
            "Certifications credited by name",
            "Suggested interview questions targeted at this candidate",
        ],
        "Explainability panel — narrative, strengths, gaps and interview questions",
        GREEN,
    )

    # 27 — search ------------------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  09",
        "Search in plain English, answered by the graph",
        "The query is parsed into skills and constraints, expanded through the graph, then fused with vector hits.",
        "asa-15-search-answer.png",
        [
            "Five modes  ·  hybrid, semantic, keyword, graph and skill",
            "A RAG-generated answer summarising the strongest candidates",
            "The graph paths used are printed with their weights",
            "This query returned in 20 ms",
        ],
        "Search — natural-language query with a RAG answer and graph citations",
        AMBER,
    )

    # 28 — knowledge graph ---------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  10",
        "The knowledge graph, explorable",
        "The asset that grows with every resume: an interactive map of who knows what, and how it all connects.",
        "asa-10-graph-canvas.png",
        [
            "356 nodes and 1 351 relationships from just five demo resumes",
            "Focus on a skill, set traversal depth, expand any node",
            "Rebuild on demand; NetworkX in-process or Neo4j for scale",
            "Powers similar-candidate discovery and skill expansion",
        ],
        "Knowledge graph explorer — force-directed view of the live graph",
        TEAL,
    )

    # 29 — reports -----------------------------------------------------------
    deck.shot_slide(
        "Walkthrough  ·  11",
        "Reports that prove what the pipeline is doing",
        "Operational KPIs, skill-gap analysis and exports for the numbers your leadership already tracks.",
        "asa-17-reports.png",
        [
            "Parse success rate, average parse time and shortlist rate",
            "Hiring trends, skill distribution and gap analysis by period",
            "Export to PDF, CSV or Excel in one click",
            "Full audit trail available to HR Admins",
        ],
        "Reports — KPIs, trends and exportable analytics",
        RED,
    )

    # 30 — impact and next steps --------------------------------------------
    slide = deck.content_slide(
        "Impact & next steps  ·  16",
        "What changes on day one, and where we go next",
        "The platform is built, tested and running. The remaining work is your data and your rollout.",
        GREEN,
    )
    before_after = [
        ("Screening a role", "12.5 hours of reading", "Minutes, over the entire pool"),
        ("Ranking", "Depends who read it", "Same inputs, same result, every time"),
        ("Justification", "A recruiter's recollection", "Score breakdown plus graph evidence"),
        ("Past applicants", "Effectively lost", "Resurface automatically for new roles"),
    ]
    deck.rect(slide, MARGIN, 2.15, 7.55, 3.05, fill=WHITE, line=BORDER, radius=0.08)
    frame = deck.text(slide, MARGIN + 0.3, 2.35, 3.2, 0.28)
    deck.line(frame, "TODAY", 9.5, RED, bold=True, first=True)
    frame = deck.text(slide, MARGIN + 4.35, 2.35, 3.0, 0.28)
    deck.line(frame, "WITH AI SKILL ANALYSER", 9.5, GREEN, bold=True, first=True)
    for index, (label, before, after) in enumerate(before_after):
        y = 2.72 + index * 0.6
        frame = deck.text(slide, MARGIN + 0.3, y, 1.9, 0.25)
        deck.line(frame, label, 10, MIST, first=True)
        frame = deck.text(slide, MARGIN + 0.3, y + 0.21, 3.9, 0.28)
        deck.line(frame, before, 11, SLATE, first=True)
        frame = deck.text(slide, MARGIN + 4.35, y + 0.21, 3.1, 0.28)
        deck.line(frame, after, 11, INK, bold=True, first=True)
    deck.rect(slide, 8.6, 2.15, 3.98, 3.05, fill=NAVY, radius=0.08)
    frame = deck.text(slide, 8.88, 2.38, 3.4, 0.28)
    deck.line(frame, "PROPOSED NEXT STEPS", 9.5, RGBColor(0x9F, 0xC4, 0xFF), bold=True, first=True)
    steps = [
        ("Week 1", "Load your Skills CSV and a sample of 200 real resumes"),
        ("Week 2", "Calibrate weights with your recruiters on two live roles"),
        ("Week 3", "Deploy to your cloud, connect SSO, import history"),
        ("Week 4", "Pilot with one hiring team, measure time-to-shortlist"),
    ]
    for index, (when, what) in enumerate(steps):
        y = 2.72 + index * 0.62
        frame = deck.text(slide, 8.88, y, 3.4, 0.24)
        deck.line(frame, when, 10, VIOLET, bold=True, first=True)
        frame = deck.text(slide, 8.88, y + 0.2, 3.4, 0.34)
        deck.line(frame, what, 10, WHITE, first=True, spacing=1.15)
    deck.rect(slide, MARGIN, 5.42, BODY_W, 1.05, fill=RGBColor(0xEC, 0xFA, 0xF2), line=RGBColor(0xB8, 0xEC, 0xD2), radius=0.08)
    frame = deck.text(slide, MARGIN + 0.42, 5.62, 8.4, 0.6)
    deck.line(
        frame,
        "The platform is complete and demonstrable today — 138 automated tests, "
        "Docker deployment, and full documentation.",
        13,
        INK,
        bold=True,
        first=True,
        spacing=1.15,
    )
    deck.chip(slide, 9.55, 5.72, 3.03, "Ready for your data", GREEN, h=0.44, size=12)
    deck.notes(
        slide,
        "Close by asking for two things: their Skills CSV and 200 anonymised resumes. That is everything needed to run "
        "the same demo on their own data next week.",
    )

    deck.save()
    print(f"{OUTPUT}  ·  {deck.count} slides  ·  {OUTPUT.stat().st_size // 1024} KB")


if __name__ == "__main__":
    build()
